import argparse
import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from fastapi.testclient import TestClient

from app import app as fastapi_app


def extract_first_picture(pptx_path: str, slide_idx: int = 12) -> str:
    """Extrai a primeira figura (MSO_SHAPE_TYPE.PICTURE) do slide.

    Útil para inspecionar rasterização/arquivos de imagem antes de enviar ao /convert.
    """
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ext = shape.image.ext or "png"
            out_name = f"test_{shape.name}.{ext}"
            with open(out_name, "wb") as f:
                f.write(shape.image.blob)
            return out_name

    raise RuntimeError(f"Nenhuma imagem encontrada no slide_idx={slide_idx} do arquivo {pptx_path}")


def post_convert(pptx_path: str) -> dict:
    """Faz POST /convert (multipart) usando TestClient.

    Isso chama o endpoint "in-process", então breakpoints no `app.py`
    devem bater diretamente no fluxo da função.
    """
    client = TestClient(fastapi_app)

    pptx_abs = str(Path(pptx_path).resolve())
    filename = os.path.basename(pptx_abs)
    mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    with open(pptx_abs, "rb") as f:
        response = client.post(
            "/convert",
            files={"file": (filename, f, mime_type)},
            headers={"Accept": "application/json"},
        )

    try:
        payload = response.json()
    except Exception as e:
        raise RuntimeError(f"Resposta nao-JSON. status={response.status_code}.") from e

    if response.status_code >= 400:
        raise RuntimeError(f"Falha no /convert. status={response.status_code}. payload={payload}")

    return payload


def main():
    parser = argparse.ArgumentParser(description="Debug local chamando POST /convert.")
    parser.add_argument(
        "--pptx",
        default="data-raw/2025-12-22_cofin_final.pptx",
        help="Caminho do .pptx para enviar ao /convert.",
    )
    parser.add_argument(
        "--out-md",
        default="debug_convert_output.md",
        help="Arquivo para salvar o markdown retornado.",
    )
    parser.add_argument(
        "--extract-first-picture",
        action="store_true",
        help="Extrai a primeira imagem de um slide antes de chamar /convert (somente para inspeção local).",
    )
    parser.add_argument(
        "--slide-idx",
        type=int,
        default=12,
        help="Índice zero-based do slide para extração da imagem (quando --extract-first-picture).",
    )

    args = parser.parse_args()

    if args.extract_first_picture:
        extracted = extract_first_picture(args.pptx, slide_idx=args.slide_idx)
        print(f"[DEBUG] Imagem extraída: {extracted}")

    response = post_convert(args.pptx)
    markdown = response.get("markdown", "")

    print(f"[DEBUG] JSON keys: {list(response.keys())}")
    print(f"[DEBUG] markdown size: {len(markdown)} chars")
    print("[DEBUG] markdown preview:")
    print(markdown[:800])

    Path(args.out_md).write_text(markdown, encoding="utf-8")
    print(f"[DEBUG] Salvo em: {args.out_md}")


if __name__ == "__main__":
    main()