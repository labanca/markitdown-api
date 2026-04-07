import logging
import os
import re
import shutil
import subprocess
import tempfile
import time
from uuid import uuid4

from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, HTTPException

load_dotenv()

_log_level = os.getenv("LOG_LEVEL", "INFO").upper()
logging.basicConfig(
    level=_log_level,
    format="%(levelname)s %(name)s: %(message)s",
)
logger = logging.getLogger(__name__)

# Suppress verbose Azure SDK HTTP tracing unless the app itself is in DEBUG mode.
# At INFO the SDK logs every request/response at INFO level, which is noisy in
# normal operation. Keep them visible only when the developer explicitly asks
# for DEBUG so they still appear during troubleshooting.
_azure_log_level = logging.DEBUG if _log_level == "DEBUG" else logging.WARNING
logging.getLogger("azure").setLevel(_azure_log_level)

from markitdown import MarkItDown
from azure.ai.contentunderstanding import ContentUnderstandingClient
from azure.ai.contentunderstanding.models import AnalysisInput
from azure.core.credentials import AzureKeyCredential
from azure.core.exceptions import HttpResponseError

app = FastAPI(title="MarkItDown API - Conversão para Markdown")

# ---------------------------------------------------------------------------
# Upload size limit
# ---------------------------------------------------------------------------

# Maximum accepted upload size in bytes. Configurable via MAX_FILE_SIZE_MB env var.
# Defaults to 100 MB. The check is done after reading the payload so oversized
# uploads are rejected before any conversion work begins.
MAX_FILE_SIZE = int(os.getenv("MAX_FILE_SIZE_MB", "100")) * 1024 * 1024


# ---------------------------------------------------------------------------
# Azure Content Understanding client
# ---------------------------------------------------------------------------

def get_cu_client() -> ContentUnderstandingClient:
    endpoint = os.getenv("AZURE_CONTENT_UNDERSTANDING_ENDPOINT")
    key = os.getenv("AZURE_FOUNDRY_API_KEY")
    return ContentUnderstandingClient(
        endpoint=endpoint,
        credential=AzureKeyCredential(key),
        api_version="2025-11-01",
    )


# ---------------------------------------------------------------------------
# Vector image rasterization: EMF/WMF -> PNG via LibreOffice
# ---------------------------------------------------------------------------

_VECTOR_FORMATS = {"wmf", "emf", "svg"}

_EXT_TO_MIME = {
    "jpg": "image/jpeg", "jpeg": "image/jpeg",
    "png": "image/png", "bmp": "image/bmp",
    "tif": "image/tiff", "tiff": "image/tiff",
}


def _resolve_libreoffice_bin() -> str:
    """Resolve o executavel do LibreOffice em Linux/Windows.

    Prioridade:
    1) Variavel de ambiente LIBREOFFICE_BIN.
    2) PATH (libreoffice, soffice, soffice.exe).
    """
    env_bin = os.getenv("LIBREOFFICE_BIN")
    if env_bin:
        return env_bin

    for candidate in ("libreoffice", "soffice", "soffice.exe"):
        path = shutil.which(candidate)
        if path:
            return path

    raise RuntimeError(
        "LibreOffice nao encontrado no PATH. Defina LIBREOFFICE_BIN "
        "(ex: C:/Program Files/LibreOffice/program/soffice.exe) "
        "ou instale/adicione o LibreOffice ao PATH."
    )


def convert_to_pdf(image_bytes: bytes, ext: str) -> bytes:
    """Convert a vector image (EMF, WMF, SVG) to PDF using LibreOffice headless.

    LibreOffice preserves vector information when converting to PDF, so the
    Content Understanding service receives a high-fidelity document rather than
    a rasterized bitmap — no resolution loss, no interpolation artifacts.
    """
    with tempfile.TemporaryDirectory() as tmp_dir:
        src = os.path.join(tmp_dir, f"input.{ext}")
        with open(src, "wb") as f:
            f.write(image_bytes)

        libreoffice_bin = _resolve_libreoffice_bin()

        result = subprocess.run(
            [
                libreoffice_bin, "--headless", "--norestore",
                "--convert-to", "pdf",
                "--outdir", tmp_dir,
                src,
            ],
            capture_output=True,
            timeout=30,
        )

        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice PDF conversion failed: {result.stderr.decode(errors='replace')}"
            )

        dst = os.path.join(tmp_dir, "input.pdf")
        if not os.path.exists(dst):
            raise RuntimeError(
                f"LibreOffice produced no PDF output. "
                f"stdout: {result.stdout.decode(errors='replace')}"
            )

        with open(dst, "rb") as f:
            return f.read()


# ---------------------------------------------------------------------------
# Table extraction via Azure Content Understanding Layout
# ---------------------------------------------------------------------------

def _parse_source_midpoint(source: str) -> tuple[float, float] | None:
    """Extract (x_mid, y_mid) from an Azure CU source string.

    Format: D(page, x1,y1, x2,y2, x3,y3, x4,y4)
    where the four corners are topLeft, topRight, bottomRight, bottomLeft.
    """
    nums = [float(n) for n in re.findall(r"[\d.]+", source)]
    if len(nums) < 9:
        return None
    x_mid = (nums[1] + nums[3]) / 2          # average of topLeft.x and topRight.x
    y_mid = (nums[2] + nums[8]) / 2          # average of topLeft.y and bottomLeft.y
    return x_mid, y_mid


def _reconstruct_table_from_lines(lines: list) -> str | None:
    """Reconstruct a markdown table from spatially-positioned lines.

    Used as a fallback when Azure Content Understanding returns content as
    paragraphs instead of a structured table — which happens with borderless
    tables (no visible grid lines).

    Algorithm:
    1. Parse each line's bounding box from the 'source' coordinate string.
    2. Cluster lines into rows by Y-midpoint proximity.
    3. Within each row, sort lines by X-midpoint (left → right).
    4. Use the widest row as the column reference to assign cells.
    5. Emit a Markdown table.

    Returns None if the content does not look like a multi-column table.
    """
    items: list[tuple[float, float, str]] = []
    for line in lines:
        source = line.get("source", "")
        content = line.get("content", "").replace("\n", " ").replace("|", "\\|").strip()
        if not source or not content:
            continue
        coords = _parse_source_midpoint(source)
        if coords is None:
            continue
        x_mid, y_mid = coords
        items.append((y_mid, x_mid, content))

    if not items:
        return None

    items.sort(key=lambda t: (t[0], t[1]))

    # Cluster into rows: lines whose Y-midpoints are within tolerance belong together
    y_tolerance = 8.0
    rows: list[list[tuple[float, float, str]]] = []
    current_row: list[tuple[float, float, str]] = [items[0]]
    for item in items[1:]:
        if abs(item[0] - current_row[0][0]) <= y_tolerance:
            current_row.append(item)
        else:
            rows.append(sorted(current_row, key=lambda t: t[1]))
            current_row = [item]
    rows.append(sorted(current_row, key=lambda t: t[1]))

    # Only reconstruct if it looks like a table: >= 2 rows each with >= 2 items
    if sum(1 for r in rows if len(r) >= 2) < 2:
        return None

    # Use the row with the most items as the column reference (typically the header)
    ref_row = max(rows, key=len)
    n_cols = len(ref_row)
    col_centers = [item[1] for item in ref_row]

    def nearest_col(x: float) -> int:
        return min(range(n_cols), key=lambda i: abs(col_centers[i] - x))

    grid_rows = []
    for row in rows:
        grid = [""] * n_cols
        for _, x_mid, content in row:
            grid[nearest_col(x_mid)] = content
        grid_rows.append(grid)

    md_lines = ["| " + " | ".join(grid_rows[0]) + " |",
                "| " + " | ".join(["---"] * n_cols) + " |"]
    for row in grid_rows[1:]:
        md_lines.append("| " + " | ".join(row) + " |")

    return "\n".join(md_lines)


def _tables_to_markdown(tables: list) -> str:
    """Convert Content Understanding table objects to markdown tables.

    Uses the structured cells array (rowIndex, columnIndex, content) to build
    markdown directly — no HTML parsing needed.
    """
    md_tables = []
    for table in tables:
        cells = table.get("cells", [])
        row_count = table.get("rowCount", 0)
        col_count = table.get("columnCount", 0)

        if not cells or row_count == 0 or col_count == 0:
            logger.warning(
                "Skipping table with unusable structure "
                "(cells=%d, rowCount=%d, columnCount=%d)",
                len(cells), row_count, col_count,
            )
            continue

        # Build a 2D grid filled with empty strings
        grid = [[""] * col_count for _ in range(row_count)]
        header_row = set()

        for cell in cells:
            r = cell.get("rowIndex", 0)
            c = cell.get("columnIndex", 0)
            content = (
                cell.get("content", "")
                .replace("\n", " ")
                .replace("|", "\\|")   # escape pipes to keep Markdown table valid
                .strip()
            )
            if r < row_count and c < col_count:
                grid[r][c] = content
            if cell.get("kind") == "columnHeader":
                header_row.add(r)

        lines = []
        for r_idx, row in enumerate(grid):
            line = "| " + " | ".join(row) + " |"
            lines.append(line)
            # Add separator after the last header row
            if r_idx in header_row and (r_idx + 1) not in header_row:
                sep = "| " + " | ".join(["---"] * col_count) + " |"
                lines.append(sep)

        md_tables.append("\n".join(lines))

    return "\n\n".join(md_tables)


def extract_tables_from_image(
    image_bytes: bytes,
    mime_type: str,
    image_name: str,
) -> str | None:
    """Send image bytes to Azure Content Understanding Layout and return markdown.

    Returns a markdown string with the extracted tables, or None if no tables
    were found. Raises on API errors (caller handles retry if needed).
    """
    client = get_cu_client()

    poller = client.begin_analyze(
        analyzer_id="prebuilt-layout",
        inputs=[AnalysisInput(data=image_bytes, mime_type=mime_type)],
    )
    result = poller.result()
    result_dict = result.as_dict()

    inner = result_dict.get("result", result_dict)  # handle both wrapped and unwrapped
    contents = inner.get("contents", [])
    if not contents:
        return None

    content = contents[0]
    tables = content.get("tables", [])

    if tables:
        return _tables_to_markdown(tables)

    # Fallback: Azure CU did not detect a formal table structure.
    # This happens with borderless tables (no visible grid lines) or tables
    # with thin grid lines that the layout model does not detect as cells.
    # Attempt to reconstruct the table from the spatial positions of the
    # text elements using their bounding-box coordinates.
    #
    # Azure CU response structure varies by image/API version:
    #   - Some responses expose 'lines'/'paragraphs' at contents[0] level
    #   - Others nest them inside contents[0].pages[0]
    # We probe both locations, preferring 'lines' (more granular) over
    # 'paragraphs', and the page level over the content level when both exist.
    _page0 = next(iter(content.get("pages") or []), {})
    cu_lines = (
        _page0.get("lines")
        or _page0.get("paragraphs")
        or content.get("lines")
        or content.get("paragraphs", [])
    )
    if cu_lines:
        reconstructed = _reconstruct_table_from_lines(cu_lines)
        if reconstructed:
            logger.info("Reconstructed borderless table from spatial lines in %s", image_name)
            return reconstructed

    return None


# ---------------------------------------------------------------------------
# Per-image processing: detect format, rasterize if needed, extract tables
# ---------------------------------------------------------------------------

# HTTP status codes that indicate a permanent failure — retrying will not help.
_PERMANENT_HTTP_ERRORS = {400, 401, 403, 404}


def process_image(
    image_bytes: bytes,
    ext: str,
    image_name: str,
    max_retries: int = 3,
) -> str | None:
    """Process a single image extracted from a PPTX slide.

    - Vector formats (EMF/WMF/SVG): rasterize with LibreOffice first.
    - Raster formats (PNG/JPG/...): send directly to Azure CU.

    Retry logic:
    - Permanent HTTP errors (400/401/403/404) abort immediately — no retry.
    - Transient errors (5xx, timeouts, network) use exponential backoff: 1s, 2s.

    Returns extracted Markdown string, or None if no table was found.
    """
    ext = ext.lower()

    if ext in _VECTOR_FORMATS:
        try:
            image_bytes = convert_to_pdf(image_bytes, ext)
            ext = "pdf"
            logger.info("Converted to PDF: %s", image_name)
        except Exception as exc:
            logger.error("Could not rasterize %s: %s", image_name, exc)
            return None  # treated as "no table" — placeholder removed cleanly

    mime_type = "application/pdf" if ext == "pdf" else _EXT_TO_MIME.get(ext, "image/png")

    for attempt in range(max_retries):
        try:
            result = extract_tables_from_image(image_bytes, mime_type, image_name)
            logger.info("%s in %s", "Tables found" if result else "No table", image_name)
            return result
        except HttpResponseError as exc:
            if exc.status_code in _PERMANENT_HTTP_ERRORS:
                logger.error(
                    "Permanent Azure error for %s (HTTP %s) — aborting retries",
                    image_name, exc.status_code,
                )
                break
            logger.error(
                "Transient Azure error for %s (attempt %d/%d, HTTP %s, code=%s)",
                image_name, attempt + 1, max_retries, exc.status_code, exc.error_code,
            )
        except Exception as exc:
            logger.error(
                "Content Understanding failed for %s (attempt %d/%d): %s",
                image_name, attempt + 1, max_retries, type(exc).__name__,
            )

        if attempt < max_retries - 1:
            time.sleep(2 ** attempt)  # backoff: 1s before attempt 2, 2s before attempt 3

    logger.warning("All retries exhausted for %s — skipping image", image_name)
    return None


# ---------------------------------------------------------------------------
# PPTX processing: MarkItDown text + per-image Content Understanding
# ---------------------------------------------------------------------------

def _extract_slide_images(pptx_path: str) -> list[list[tuple[str, bytes, str]]]:
    """Return a list (one per slide) of [(shape_name, image_bytes, ext), ...].

    Uses python-pptx directly to avoid shape name collisions across slides.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    result = []
    for slide in prs.slides:
        slide_images: list[tuple[str, bytes, str]] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    slide_images.append(
                        (shape.name, shape.image.blob, shape.image.ext or "png")
                    )
                except Exception as exc:
                    logger.warning("Could not read image from shape '%s': %s", shape.name, exc)
        result.append(slide_images)
    return result


_SLIDE_HEADER = re.compile(r"(<!-- Slide number: (\d+) -->)")
_IMG_PLACEHOLDER = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")


def process_pptx(file_path: str) -> str:
    """Process a PPTX file and return clean markdown with tables extracted from images.

    Strategy:
    1. MarkItDown without LLM -> reliable text/structure extraction.
    2. python-pptx -> extract image bytes per slide (avoids shape name collisions).
    3. For each empty image placeholder:
       - Rasterize if EMF/WMF (LibreOffice).
       - Send to Azure Content Understanding Layout.
       - Replace placeholder with structured markdown table.
    """
    # Step 1: base markdown (text only, no LLM)
    base_markdown = MarkItDown().convert(file_path).text_content

    # Step 2: extract images per slide
    slide_images = _extract_slide_images(file_path)

    # Step 3: walk slide by slide and replace placeholders
    parts = _SLIDE_HEADER.split(base_markdown)
    # layout: [pre, full_header, slide_num, content, ...]

    output: list[str] = []
    i = 0

    if parts:
        output.append(parts[0])
        i = 1

    while i < len(parts):
        full_header   = parts[i]
        slide_num_str = parts[i + 1]
        content       = parts[i + 2] if i + 2 < len(parts) else ""
        i += 3

        output.append(full_header)

        slide_idx = int(slide_num_str) - 1  # 0-based
        images_for_slide = slide_images[slide_idx] if slide_idx < len(slide_images) else []
        img_counter = [0]

        def replace_placeholder(match: re.Match) -> str:
            pos = img_counter[0]
            img_counter[0] += 1

            if pos >= len(images_for_slide):
                return ""  # no image data — remove placeholder

            shape_name, img_bytes, ext = images_for_slide[pos]
            label = f"slide{slide_idx + 1}_{shape_name}.{ext}"
            logger.info("Processing image: %s", label)

            extracted = process_image(img_bytes, ext, label)

            if extracted is None:
                return ""  # no table in this image — remove placeholder cleanly
            return f"\n{extracted}\n"

        content = _IMG_PLACEHOLDER.sub(replace_placeholder, content)
        output.append(content)

    return "".join(output)


# ---------------------------------------------------------------------------
# Main /convert endpoint
# ---------------------------------------------------------------------------

@app.post("/convert")
async def convert_to_markdown(file: UploadFile):
    if not file.filename:
        raise HTTPException(status_code=400, detail="Nenhum arquivo enviado")

    ext = file.filename.rsplit(".", 1)[-1].lower()

    # Reject legacy binary PowerPoint format early with a clear message.
    # python-pptx only supports .pptx (OOXML); .ppt requires a dedicated converter.
    if ext == "ppt":
        raise HTTPException(
            status_code=400,
            detail=(
                "Formato .ppt (PowerPoint 97-2003) não é suportado. "
                "Converta para .pptx e tente novamente."
            ),
        )

    # Read the full payload first so we can enforce the size limit before
    # writing anything to disk or starting any conversion work.
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"Arquivo muito grande. Limite: {MAX_FILE_SIZE // 1024 // 1024} MB",
        )

    unique_id = uuid4().hex
    temp_dir = f"/app/temp/{unique_id}"
    os.makedirs(temp_dir, exist_ok=True)
    file_path = os.path.join(temp_dir, file.filename)

    try:
        with open(file_path, "wb") as buffer:
            buffer.write(content)

        if ext == "pptx":
            markdown_content = process_pptx(file_path)
        else:
            # All other formats: standard MarkItDown (no LLM needed for text docs)
            markdown_content = MarkItDown().convert(file_path).text_content

        return {"markdown": markdown_content}

    except HTTPException:
        raise
    except Exception:
        # Log the full traceback internally; return a generic message to the client
        # so stack traces, file paths and credentials are never exposed.
        logger.exception("Conversion failed for file '%s'", file.filename)
        raise HTTPException(status_code=500, detail="Erro interno na conversão")

    finally:
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir, ignore_errors=True)


@app.get("/health")
async def health_check():
    return {"status": "healthy"}


# ---------------------------------------------------------------------------
# Debug endpoint: convert EMF/WMF to PDF and return as download
# ---------------------------------------------------------------------------

@app.post("/debug-to-pdf")
async def debug_to_pdf(file: UploadFile):
    """Recebe um arquivo EMF/WMF e devolve o PDF gerado pelo LibreOffice.

    Util para inspecionar visualmente o que o LibreOffice esta gerando
    antes de enviar ao Content Understanding.

    Exemplo de uso:
        curl -X POST http://localhost:8000/debug-to-pdf \
             -F "file=@imagem.wmf" \
             --output resultado.pdf
    """
    from fastapi.responses import Response

    if not file.filename:
        raise HTTPException(status_code=400, detail="Nenhum arquivo enviado")

    ext = file.filename.rsplit(".", 1)[-1].lower()
    if ext not in _VECTOR_FORMATS:
        raise HTTPException(
            status_code=400,
            detail=f"Formato nao suportado: .{ext}. Envie um arquivo EMF, WMF ou SVG.",
        )

    image_bytes = await file.read()

    try:
        pdf_bytes = convert_to_pdf(image_bytes, ext)
    except Exception:
        logger.exception("PDF conversion failed for '%s'", file.filename)
        raise HTTPException(status_code=500, detail="Erro interno na conversão para PDF")

    filename = file.filename.rsplit(".", 1)[0] + ".pdf"
    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )